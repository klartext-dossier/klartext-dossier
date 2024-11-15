import logging, zipfile, tempfile, io

from pptx import Presentation
from pptx.util import Inches

from dm.exceptions import TaskException
from dm.tasks.Task import Task


class PowerPointTask(Task):

    """ The zip-to-pptx task.

        This task converts a set of images to a PPTX file.
    """   

    def tryConvertingToPPT(self):
        try:
            logging.info(f'{self.name} - Converting PNGs to PPTX')

            prs = Presentation()
            slide_layout = prs.slide_layouts[6]

            with zipfile.ZipFile(self.content.data, 'r') as zipf:                
                for file in zipf.infolist():
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(zipf.open(file), Inches(0), Inches(0), prs.slide_width)

            tmp = tempfile.NamedTemporaryFile('w+b', suffix='.pptx')           
            prs.save(tmp)            
            tmp.seek(0)

            return tmp            
        except Exception as e:
            raise TaskException(f'{self.name} - cannot write PPTX', e)

    def run(self, context):
        self.load()

        ppt = self.tryConvertingToPPT()   
        self.content.setData(ppt.read())

        self.save()
